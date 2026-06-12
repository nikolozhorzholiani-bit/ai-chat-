"""Text extraction from PDF, DOCX, PPTX, XLSX, TXT."""
import io, logging
log = logging.getLogger(__name__)


def extract_text(filename: str, data: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    try:
        if ext == "pdf":
            return _pdf(data)
        elif ext in ("doc", "docx"):
            return _docx(data)
        elif ext in ("ppt", "pptx"):
            return _pptx(data)
        elif ext in ("xls", "xlsx"):
            return _xlsx(data)
        elif ext == "txt":
            return data.decode("utf-8", errors="ignore")
        else:
            return ""
    except Exception as e:
        log.error("extract_text %s: %s", filename, e)
        return ""


def _pdf(data: bytes) -> str:
    import pdfplumber, io
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        return "\n".join(p.extract_text() or "" for p in pdf.pages).strip()


def _docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            lines.append(f"[სლაიდი {i}]\n" + "\n".join(slide_lines))
    return "\n\n".join(lines)


def _xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[ფურცელი: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("  |  ".join(cells))
    return "\n".join(lines)
